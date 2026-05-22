"""Build the entrega_final.pptx presentation for the KNN class-imbalance project.

Revision 2 — applied critical-review improvements:
- Standardised slide content to English (PT speaker notes preserved)
- Stronger, statement-style titles
- Slide 12 (rank convergence) moved to annex; main flow tightened to 16 slides
- Slide 9 results table compacted to 3 key metrics + visual rank emphasis
- Slide 13 statistics simplified — single dual-baseline table, no label repetition
- Slide 7 setup table compacted from 8 to 5 rows
- Slide 2 executive summary tightened to one headline per quadrant
- Speaker notes trimmed by ~20% across slides 4/14/15 to fit a 12-min target
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ───────────────────────────── Palette ───────────────────────────────
NAVY   = RGBColor(0x00, 0x3D, 0x7A)
BLUE   = RGBColor(0x0D, 0x6E, 0xFD)
LIGHT  = RGBColor(0xF8, 0xF9, 0xFA)
GRAY   = RGBColor(0x6C, 0x75, 0x7D)
DARK   = RGBColor(0x21, 0x29, 0x39)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x10, 0x82, 0x4D)
RED    = RGBColor(0xC0, 0x39, 0x2B)
AMBER  = RGBColor(0xD9, 0x82, 0x16)
PALE   = RGBColor(0xE7, 0xF0, 0xFA)

FONT_HEAD = "Calibri"
FONT_BODY = "Calibri"

PROJECT_ROOT = Path("/Users/duarte/Desktop/knn_classimbalance")
FIG_DIR = PROJECT_ROOT / "results" / "figures"
OUTPUT  = PROJECT_ROOT / "presentation" / "entrega_final.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK_LAYOUT = prs.slide_layouts[6]

# ───────────────────────────── Helpers ───────────────────────────────
TOTAL_SLIDES = 30  # fixed total; we'll pass page numbers explicitly

def add_slide(notes_pt: str = ""):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    if notes_pt:
        slide.notes_slide.notes_text_frame.text = notes_pt
    return slide


def add_rect(slide, x, y, w, h, fill=None, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    return shp


def add_text(slide, text, x, y, w, h,
             size=14, bold=False, italic=False,
             color=DARK, font=FONT_BODY, align=PP_ALIGN.LEFT,
             valign=MSO_ANCHOR.TOP, margin=0):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(margin); tf.margin_right = Emu(margin)
    tf.margin_top = Emu(margin); tf.margin_bottom = Emu(margin)
    tf.vertical_anchor = valign
    if isinstance(text, str):
        lines = text.split("\n"); first = True
        for ln in lines:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = ln
            r.font.name = font; r.font.size = Pt(size)
            r.font.bold = bold; r.font.italic = italic
            r.font.color.rgb = color
            first = False
    else:
        first = True
        for item in text:
            if isinstance(item, str):
                txt, opts = item, {}
            else:
                txt, opts = item
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.alignment = opts.get("align", align)
            r = p.add_run()
            r.text = txt
            r.font.name = opts.get("font", font)
            r.font.size = Pt(opts.get("size", size))
            r.font.bold = opts.get("bold", bold)
            r.font.italic = opts.get("italic", italic)
            r.font.color.rgb = opts.get("color", color)
            first = False
    return tx


def add_image(slide, path, x, y, w=None, h=None):
    if w is not None and h is not None:
        return slide.shapes.add_picture(str(path), x, y, w, h)
    if w is not None:
        return slide.shapes.add_picture(str(path), x, y, width=w)
    if h is not None:
        return slide.shapes.add_picture(str(path), x, y, height=h)
    return slide.shapes.add_picture(str(path), x, y)


def slide_header(slide, kicker, title, page_num):
    add_text(slide, kicker, Inches(0.55), Inches(0.35),
             Inches(8), Inches(0.3),
             size=11, bold=True, color=NAVY, font=FONT_HEAD)
    add_text(slide, title, Inches(0.55), Inches(0.65),
             Inches(11.5), Inches(0.75),
             size=28, bold=True, color=DARK, font=FONT_HEAD, margin=0)
    add_text(slide, f"{page_num} / {TOTAL_SLIDES}", Inches(12.3), Inches(7.1),
             Inches(1), Inches(0.3),
             size=9, color=GRAY, align=PP_ALIGN.RIGHT)
    add_text(slide, "KNNFairRank  ·  Machine Learning I  ·  CC2008",
             Inches(0.55), Inches(7.1), Inches(8), Inches(0.3),
             size=9, color=GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Bom dia. Somos o Grupo 2 de Machine Learning I e vamos apresentar o "
    "KNNFairRank — uma correcção ao KNN para problemas de class imbalance, "
    "derivada de order statistics de processos de Poisson e validada num "
    "benchmark de 49 datasets sob o protocolo de Demšar 2006.\n\n"
    "A apresentação tem 12 minutos. Começamos pelo executive summary.\n\n"
    "[≈25s]"
)
add_rect(s, Inches(0), Inches(0), Inches(4.5), SH, fill=NAVY)
add_rect(s, Inches(4.5), Inches(0), Inches(0.08), SH, fill=BLUE)
add_text(s, "UNIVERSIDADE DO PORTO",
         Inches(0.5), Inches(0.7), Inches(3.7), Inches(0.4),
         size=12, bold=True, color=WHITE, font=FONT_HEAD)
add_text(s, "Faculdade de Ciências",
         Inches(0.5), Inches(1.05), Inches(3.7), Inches(0.3),
         size=11, color=PALE)
add_text(s, "Machine Learning I",
         Inches(0.5), Inches(5.6), Inches(3.7), Inches(0.4),
         size=14, bold=True, color=WHITE, font=FONT_HEAD)
add_text(s, "CC2008  ·  2025/2026",
         Inches(0.5), Inches(5.95), Inches(3.7), Inches(0.3),
         size=11, color=PALE)
add_text(s, "Dataset Group 2",
         Inches(0.5), Inches(6.25), Inches(3.7), Inches(0.3),
         size=11, color=PALE)
add_text(s, "Binary classification under class imbalance",
         Inches(0.5), Inches(6.5), Inches(3.7), Inches(0.4),
         size=10, color=PALE, italic=True)

add_text(s, "PRACTICAL ASSIGNMENT",
         Inches(5.0), Inches(1.5), Inches(8.0), Inches(0.3),
         size=11, bold=True, color=NAVY, font=FONT_HEAD)
add_text(s, "KNN Variants for\nClass Imbalance",
         Inches(5.0), Inches(1.85), Inches(8.0), Inches(1.7),
         size=44, bold=True, color=DARK, font=FONT_HEAD)
add_text(s, "A Rank-Correction Approach Derived from\nPoisson Order Statistics",
         Inches(5.0), Inches(3.65), Inches(8.0), Inches(0.9),
         size=20, color=NAVY, italic=True, font=FONT_HEAD)

add_text(s, "AUTHORS", Inches(5.0), Inches(4.95), Inches(8.0), Inches(0.3),
         size=11, bold=True, color=GRAY, font=FONT_HEAD)
add_rect(s, Inches(5.0), Inches(5.25), Inches(1.5), Inches(0.02), fill=NAVY)
_authors = [
    ("Duarte Gomes", "up202409386"),
    ("José Sousa",   "up202405046"),
    ("Tiago Sousa",  "up202405406"),
]
for _i, (_name, _num) in enumerate(_authors):
    _y = Inches(5.45 + _i * 0.45)
    add_text(s, _name, Inches(5.0), _y, Inches(3.0), Inches(0.4),
             size=15, bold=True, color=DARK, font=FONT_HEAD,
             valign=MSO_ANCHOR.MIDDLE)
    add_text(s, _num, Inches(7.6), _y, Inches(2.5), Inches(0.4),
             size=14, color=GRAY, font=FONT_HEAD,
             valign=MSO_ANCHOR.MIDDLE)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — Executive Summary  (compacted: one headline per quadrant)
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Em 30 segundos: o problema é que o KNN clássico falha sob class "
    "imbalance por um enviesamento estrutural — não por escolha de k. "
    "A nossa proposta corrige isso usando uma comparação de rank "
    "derivada de Poisson. Validámos em 49 datasets com a battery "
    "completa de Demšar. E os resultados batem estatisticamente o "
    "KNNOptK e o SMOTE com effect sizes grandes.\n\n"
    "[≈40s. Cada quadrante vai ser retomado em profundidade nas próximas "
    "slides — esta é apenas a vista global.]"
)
slide_header(s, "EXECUTIVE SUMMARY", "The project in four ideas", 2)

quad_w = Inches(5.9); quad_h = Inches(2.55)
quad_x1, quad_x2 = Inches(0.55), Inches(6.85)
quad_y1, quad_y2 = Inches(1.7),  Inches(4.4)

def quadrant(x, y, kicker, kicker_color, headline, sub):
    add_rect(s, x, y, quad_w, quad_h, fill=LIGHT)
    add_rect(s, x, y, Inches(0.08), quad_h, fill=kicker_color)
    add_text(s, kicker, x + Inches(0.3), y + Inches(0.22),
             quad_w - Inches(0.4), Inches(0.3),
             size=11, bold=True, color=kicker_color, font=FONT_HEAD)
    add_text(s, headline, x + Inches(0.3), y + Inches(0.7),
             quad_w - Inches(0.4), Inches(1.2),
             size=22, bold=True, color=DARK, font=FONT_HEAD,
             valign=MSO_ANCHOR.MIDDLE)
    add_text(s, sub, x + Inches(0.3), y + Inches(1.85),
             quad_w - Inches(0.4), Inches(0.6),
             size=13, italic=True, color=GRAY)

quadrant(quad_x1, quad_y1, "PROBLEM", RED,
         "Structural bias",
         "KNN compares  d₁ᵐⁱⁿ  to  d₁ᵐᵃʲ  — biased by sampling alone.")
quadrant(quad_x2, quad_y1, "APPROACH", BLUE,
         "Fair rank: kₑff = r",
         "Compare  d₁ᵐⁱⁿ  to  d_rᵐᵃʲ — dimension cancels in the derivation.")
quadrant(quad_x1, quad_y2, "VALIDATION", NAVY,
         "Demšar 2006 protocol",
         "10 algorithms × 49 datasets × 50 splits  +  Instance Space Analysis.")
quadrant(quad_x2, quad_y2, "RESULTS", GREEN,
         "Statistically  +  practically",
         "Beats KNNOptK (d = +0.88) and SMOTE (d = +0.56); ρ = +0.77 theory↔data.")


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — KNN: how it works
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Recordemos o KNN. Para classificar uma query, vamos buscar os k "
    "vizinhos mais próximos e votamos por maioria. Não há treino: o "
    "algoritmo só memoriza X e y.\n\n"
    "Pontos fortes à direita: é não-paramétrico, interpretável e simples. "
    "Pontos fracos: lazy, sensível à escala — e estruturalmente enviesado "
    "sob class imbalance, que é exactamente onde o nosso trabalho actua.\n\n"
    "[≈35s. Transição: já a seguir vamos ver porquê o enviesamento é "
    "estrutural e não-removível por escolha de k.]"
)
slide_header(s, "BASE ALGORITHM", "How KNN decides", 3)

diag_x, diag_y = Inches(0.7), Inches(1.7)
diag_w, diag_h = Inches(5.5), Inches(4.7)
add_rect(s, diag_x, diag_y, diag_w, diag_h, fill=LIGHT)

maj_pts = [(1.2, 0.5), (1.7, 1.0), (2.4, 0.7), (3.1, 1.2), (4.0, 0.5),
           (4.5, 1.4), (3.8, 2.4), (4.7, 2.8), (1.1, 2.7), (0.6, 1.5),
           (0.8, 3.2), (1.9, 3.6), (3.0, 3.7), (4.4, 3.6), (2.6, 2.0)]
min_pts = [(2.2, 1.7), (2.8, 1.6), (2.5, 2.4)]

for (px, py) in maj_pts:
    cx, cy = diag_x + Inches(px - 0.13), diag_y + Inches(py - 0.13)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, Inches(0.26), Inches(0.26))
    dot.fill.solid(); dot.fill.fore_color.rgb = BLUE
    dot.line.color.rgb = WHITE; dot.line.width = Pt(1.5)
    dot.shadow.inherit = False
for (px, py) in min_pts:
    cx, cy = diag_x + Inches(px - 0.14), diag_y + Inches(py - 0.14)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, Inches(0.28), Inches(0.28))
    dot.fill.solid(); dot.fill.fore_color.rgb = RED
    dot.line.color.rgb = WHITE; dot.line.width = Pt(1.5)
    dot.shadow.inherit = False

qx, qy = 2.5, 2.0
cx, cy = diag_x + Inches(qx - 0.20), diag_y + Inches(qy - 0.20)
qshape = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, Inches(0.40), Inches(0.40))
qshape.fill.solid(); qshape.fill.fore_color.rgb = NAVY
qshape.line.color.rgb = WHITE; qshape.line.width = Pt(2.5)
qshape.shadow.inherit = False

ring_r = 1.2
ring = s.shapes.add_shape(MSO_SHAPE.OVAL,
                          diag_x + Inches(qx - ring_r),
                          diag_y + Inches(qy - ring_r),
                          Inches(ring_r * 2), Inches(ring_r * 2))
ring.fill.background()
ring.line.color.rgb = NAVY; ring.line.width = Pt(2); ring.line.dash_style = 7
ring.shadow.inherit = False

add_text(s, "Query  +  k nearest neighbours",
         diag_x + Inches(0.3), diag_y + Inches(4.05),
         diag_w - Inches(0.6), Inches(0.3),
         size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font=FONT_HEAD)
add_text(s, "Majority vote → predicted class",
         diag_x + Inches(0.3), diag_y + Inches(4.32),
         diag_w - Inches(0.6), Inches(0.3),
         size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

leg_y = diag_y + Inches(0.2)
def legend_dot(x, color):
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, x, leg_y, Inches(0.18), Inches(0.18))
    d.fill.solid(); d.fill.fore_color.rgb = color
    d.line.fill.background(); d.shadow.inherit = False
legend_dot(diag_x + Inches(0.25), BLUE)
add_text(s, "majority", diag_x + Inches(0.5), leg_y - Inches(0.04),
         Inches(1.2), Inches(0.25), size=10, color=DARK)
legend_dot(diag_x + Inches(1.7), RED)
add_text(s, "minority", diag_x + Inches(1.95), leg_y - Inches(0.04),
         Inches(1.2), Inches(0.25), size=10, color=DARK)
legend_dot(diag_x + Inches(3.2), NAVY)
add_text(s, "query", diag_x + Inches(3.45), leg_y - Inches(0.04),
         Inches(1.2), Inches(0.25), size=10, color=DARK)

card_x = Inches(6.8); card_w = Inches(6.05)
add_rect(s, card_x, Inches(1.7), card_w, Inches(2.2), fill=LIGHT)
add_rect(s, card_x, Inches(1.7), Inches(0.08), Inches(2.2), fill=GREEN)
add_text(s, "STRENGTHS", card_x + Inches(0.25), Inches(1.85),
         card_w, Inches(0.3),
         size=11, bold=True, color=GREEN, font=FONT_HEAD)
add_text(s, [
    ("Non-parametric — no distribution assumed", {"bold": True, "size": 14}),
    ("Fully interpretable — see which neighbours voted", {"size": 14}),
    ("No training phase — just stores X, y", {"size": 14}),
], card_x + Inches(0.25), Inches(2.15), card_w - Inches(0.4), Inches(1.7),
   size=14, color=DARK)

add_rect(s, card_x, Inches(4.15), card_w, Inches(2.2), fill=LIGHT)
add_rect(s, card_x, Inches(4.15), Inches(0.08), Inches(2.2), fill=RED)
add_text(s, "LIMITATIONS", card_x + Inches(0.25), Inches(4.3),
         card_w, Inches(0.3),
         size=11, bold=True, color=RED, font=FONT_HEAD)
add_text(s, [
    ("Lazy — O(N) per query at inference", {"bold": True, "size": 14}),
    ("Distance-sensitive — needs standardisation", {"size": 14}),
    ("Structurally biased under class imbalance", {"size": 14, "bold": True, "color": RED}),
], card_x + Inches(0.25), Inches(4.6), card_w - Inches(0.4), Inches(1.7),
   size=14, color=DARK)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — Class imbalance: structural bias  (notes trimmed)
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Aqui está o problema. Mesmo com distribuições idênticas, mais "
    "pontos de uma classe que da outra cria um enviesamento na "
    "distância esperada.\n\n"
    "Sob Poisson homogéneo, E[d_k^c] escala com (k/λ_c)^(1/d). A razão "
    "cross-class no rank 1 é r^(1/d) — onde r é N_maj/N_min. A classe "
    "maioritária está sempre mais próxima, por amostragem.\n\n"
    "Pontos importantes a sublinhar: a dimensão d aparece no expoente; "
    "e o enviesamento NÃO se resolve mudando k, porque qualquer k afecta "
    "as duas classes do mesmo modo.\n\n"
    "[≈50s. Transição: e se invertermos esta equação para derivar a "
    "correcção?]"
)
slide_header(s, "CLASS IMBALANCE", "The bias is structural, not k-related", 4)

viz_x, viz_y = Inches(0.55), Inches(1.7)
viz_w, viz_h = Inches(6.5), Inches(4.8)
add_rect(s, viz_x, viz_y, viz_w, viz_h, fill=LIGHT)

import random, math
random.seed(42)
maj_cx, maj_cy = 1.5, 2.6
for _ in range(45):
    r_ = random.gauss(0, 0.65); ang = random.uniform(0, 6.28)
    px = maj_cx + r_ * math.cos(ang); py = maj_cy + r_ * math.sin(ang)
    if 0.2 < px < 5.8 and 0.4 < py < 4.0:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  viz_x + Inches(px - 0.08),
                                  viz_y + Inches(py - 0.08),
                                  Inches(0.16), Inches(0.16))
        dot.fill.solid(); dot.fill.fore_color.rgb = BLUE
        dot.line.fill.background(); dot.shadow.inherit = False
min_cx, min_cy = 4.6, 2.6
for _ in range(6):
    r_ = random.gauss(0, 0.55); ang = random.uniform(0, 6.28)
    px = min_cx + r_ * math.cos(ang); py = min_cy + r_ * math.sin(ang)
    if 0.2 < px < 5.8 and 0.4 < py < 4.0:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  viz_x + Inches(px - 0.11),
                                  viz_y + Inches(py - 0.11),
                                  Inches(0.22), Inches(0.22))
        dot.fill.solid(); dot.fill.fore_color.rgb = RED
        dot.line.color.rgb = WHITE; dot.line.width = Pt(1.2)
        dot.shadow.inherit = False

qx, qy = 3.05, 2.5
q = s.shapes.add_shape(MSO_SHAPE.OVAL,
                       viz_x + Inches(qx - 0.18),
                       viz_y + Inches(qy - 0.18),
                       Inches(0.36), Inches(0.36))
q.fill.solid(); q.fill.fore_color.rgb = NAVY
q.line.color.rgb = WHITE; q.line.width = Pt(2.5)
q.shadow.inherit = False

line_maj = s.shapes.add_connector(1,
    viz_x + Inches(qx), viz_y + Inches(qy),
    viz_x + Inches(2.10), viz_y + Inches(2.30))
line_maj.line.color.rgb = BLUE; line_maj.line.width = Pt(3)

line_min = s.shapes.add_connector(1,
    viz_x + Inches(qx), viz_y + Inches(qy),
    viz_x + Inches(4.30), viz_y + Inches(2.20))
line_min.line.color.rgb = RED; line_min.line.width = Pt(3)

add_rect(s, viz_x + Inches(2.10), viz_y + Inches(1.65),
         Inches(0.85), Inches(0.35), fill=WHITE, line=BLUE)
add_text(s, "d₁ᵐᵃʲ", viz_x + Inches(2.10), viz_y + Inches(1.65),
         Inches(0.85), Inches(0.35),
         size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
add_rect(s, viz_x + Inches(3.55), viz_y + Inches(1.55),
         Inches(0.85), Inches(0.35), fill=WHITE, line=RED)
add_text(s, "d₁ᵐⁱⁿ", viz_x + Inches(3.55), viz_y + Inches(1.55),
         Inches(0.85), Inches(0.35),
         size=13, bold=True, color=RED, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)

add_text(s, "N_maj = 90", viz_x + Inches(0.4), viz_y + Inches(0.1),
         Inches(2.0), Inches(0.3),
         size=11, bold=True, color=BLUE, font=FONT_HEAD)
add_text(s, "N_min = 10", viz_x + Inches(4.0), viz_y + Inches(0.1),
         Inches(2.0), Inches(0.3),
         size=11, bold=True, color=RED, font=FONT_HEAD)
add_text(s, "Same underlying density λ;  different N  →  d₁ᵐⁱⁿ farther by r^(1/d) on average",
         viz_x + Inches(0.3), viz_y + Inches(4.45),
         viz_w - Inches(0.6), Inches(0.3),
         size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

right_x = Inches(7.4); right_w = Inches(5.5)
add_text(s, "Poisson order statistics",
         right_x, Inches(1.75), right_w, Inches(0.35),
         size=11, bold=True, color=NAVY, font=FONT_HEAD)
add_rect(s, right_x, Inches(2.1), right_w, Inches(1.0), fill=PALE)
add_text(s, "E[d_kᶜ]  ∝  (k / λ_c)^(1/d)",
         right_x, Inches(2.4), right_w, Inches(0.5),
         size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font=FONT_HEAD)
add_text(s, "λ_c = local density of class c    ·    d = dimension",
         right_x, Inches(2.78), right_w, Inches(0.3),
         size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, "Apply at rank 1  →  cross-class ratio",
         right_x, Inches(3.4), right_w, Inches(0.35),
         size=11, bold=True, color=NAVY, font=FONT_HEAD)
add_rect(s, right_x, Inches(3.75), right_w, Inches(1.0), fill=PALE)
add_text(s, "E[d₁ᵐⁱⁿ] / E[d₁ᵐᵃʲ]  =  r^(1/d)",
         right_x, Inches(4.05), right_w, Inches(0.5),
         size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font=FONT_HEAD)
add_text(s, "r = N_maj / N_min    —    biased by sampling alone",
         right_x, Inches(4.43), right_w, Inches(0.3),
         size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
add_rect(s, right_x, Inches(5.1), right_w, Inches(1.2), fill=NAVY)
add_text(s, "Tuning k cannot remove it.",
         right_x + Inches(0.3), Inches(5.3),
         right_w - Inches(0.4), Inches(0.45),
         size=15, bold=True, color=WHITE, font=FONT_HEAD)
add_text(s, "Every k affects both classes the same way — the bias survives.",
         right_x + Inches(0.3), Inches(5.75),
         right_w - Inches(0.4), Inches(0.6),
         size=12, color=PALE, italic=True)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — Theoretical derivation
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Aqui está a contribuição. Se E[d_k^c] depende de (k/λ_c)^(1/d), "
    "podemos encontrar um k_eff tal que E[d_{k_eff}^maj] = E[d_1^min]. "
    "Resolvendo a igualdade e elevando ambos os lados à potência d, "
    "a dimensão CANCELA. Ficamos com k_eff = r.\n\n"
    "Três propriedades a sublinhar: é dimension-free — não precisamos de "
    "saber d; é fechado em forma; e é interpretável — r é uma grandeza "
    "directa do dataset.\n\n"
    "[≈45s. Transição: vejamos a comparação com o KNN clássico.]"
)
slide_header(s, "PROPOSED MODIFICATION", "Poisson-derived fair rank", 5)

box_y, box_h = Inches(2.1), Inches(2.4)
b1_x, b2_x, b3_x = Inches(0.7), Inches(5.0), Inches(9.3)
b_w = Inches(3.4)

def deriv_box(x, kicker, formula, caption):
    add_rect(s, x, box_y, b_w, box_h, fill=LIGHT)
    add_rect(s, x, box_y, Inches(0.08), box_h, fill=NAVY)
    add_text(s, kicker, x + Inches(0.25), box_y + Inches(0.2),
             b_w, Inches(0.3),
             size=11, bold=True, color=NAVY, font=FONT_HEAD)
    add_text(s, formula, x + Inches(0.25), box_y + Inches(0.65),
             b_w - Inches(0.4), Inches(1.1),
             size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
    add_text(s, caption, x + Inches(0.25), box_y + Inches(1.75),
             b_w - Inches(0.4), Inches(0.65),
             size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

deriv_box(b1_x, "①  POISSON SCALING",
          "E[d_kᶜ] ∝ (k / λc)^(1/d)",
          "Expected distance to the k-th\nnearest neighbour of class c")

arr1 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                          Inches(4.45), Inches(3.05),
                          Inches(0.5), Inches(0.5))
arr1.fill.solid(); arr1.fill.fore_color.rgb = NAVY
arr1.line.fill.background(); arr1.shadow.inherit = False

deriv_box(b2_x, "②  EQUATE EXPECTATIONS",
          "E[d_{k_eff}ᵐᵃʲ] = E[d₁ᵐⁱⁿ]",
          "Find k_eff that makes the\nrank-1 comparison fair")

arr2 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                          Inches(8.75), Inches(3.05),
                          Inches(0.5), Inches(0.5))
arr2.fill.solid(); arr2.fill.fore_color.rgb = NAVY
arr2.line.fill.background(); arr2.shadow.inherit = False

deriv_box(b3_x, "③  DIMENSION CANCELS",
          "k_eff = r = N_maj / N_min",
          "No hyperparameter, no dimension\nestimation — closed form")

add_rect(s, Inches(0.55), Inches(5.0), Inches(12.2), Inches(1.45), fill=NAVY)
add_text(s, "Compare  d₁ᵐⁱⁿ   vs   d_rᵐᵃʲ",
         Inches(0.7), Inches(5.1), Inches(11.9), Inches(0.7),
         size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         font=FONT_HEAD, valign=MSO_ANCHOR.MIDDLE)
add_text(s, "Instead of  d₁ᵐⁱⁿ  vs  d₁ᵐᵃʲ   —   accounts for sampling-induced density imbalance",
         Inches(0.7), Inches(5.85), Inches(11.9), Inches(0.5),
         size=14, color=PALE, align=PP_ALIGN.CENTER, italic=True)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — Standard vs FairRank + variants
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Comparação directa. Standard KNN à esquerda — pega no d_1 minority "
    "e no d_1 majority. FairRank à direita — pega no d_1 minority e no "
    "d_r majority.\n\n"
    "Sobre essa base desenvolvemos uma família de 5 variantes. O core "
    "tem k_eff = r exacto. O CV ajusta α em k_eff = r^α via inner CV — "
    "vamos ver porquê na slide do Poisson empírico. O JointCV tuna α e "
    "n_votes em conjunto. O Ensemble varre o grid completo sem CV. E o "
    "Jackknife usa LOO sobre os minority ranks para reduzir variância.\n\n"
    "[≈45s. Transição: vejamos o setup experimental.]"
)
slide_header(s, "PROPOSED MODIFICATION", "Standard vs FairRank  +  family of variants", 6)

col_y, col_h = Inches(1.8), Inches(3.5)
col_l_x, col_r_x = Inches(0.6), Inches(6.95)
col_w = Inches(5.85)

add_rect(s, col_l_x, col_y, col_w, col_h, fill=LIGHT)
add_rect(s, col_l_x, col_y, col_w, Inches(0.55), fill=GRAY)
add_text(s, "STANDARD KNN", col_l_x, col_y, col_w, Inches(0.55),
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
add_text(s, "Comparison rule", col_l_x + Inches(0.3), col_y + Inches(0.8),
         col_w, Inches(0.3), size=11, bold=True, color=GRAY, font=FONT_HEAD)
add_rect(s, col_l_x + Inches(0.3), col_y + Inches(1.15),
         col_w - Inches(0.6), Inches(0.9), fill=WHITE)
add_text(s, "d₁ᵐⁱⁿ   <   d₁ᵐᵃʲ   ?",
         col_l_x + Inches(0.3), col_y + Inches(1.15),
         col_w - Inches(0.6), Inches(0.9),
         size=22, bold=True, color=DARK, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
add_text(s, "Both ranks fixed at 1.",
         col_l_x + Inches(0.3), col_y + Inches(2.2),
         col_w - Inches(0.6), Inches(0.3),
         size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
add_rect(s, col_l_x + Inches(0.3), col_y + Inches(2.65),
         col_w - Inches(0.6), Inches(0.7), fill=PALE)
add_text(s, "→  biased toward majority class",
         col_l_x + Inches(0.3), col_y + Inches(2.65),
         col_w - Inches(0.6), Inches(0.7),
         size=14, bold=True, color=RED, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)

add_rect(s, col_r_x, col_y, col_w, col_h, fill=LIGHT)
add_rect(s, col_r_x, col_y, col_w, Inches(0.55), fill=NAVY)
add_text(s, "KNNFAIRRANK", col_r_x, col_y, col_w, Inches(0.55),
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
add_text(s, "Comparison rule", col_r_x + Inches(0.3), col_y + Inches(0.8),
         col_w, Inches(0.3), size=11, bold=True, color=NAVY, font=FONT_HEAD)
add_rect(s, col_r_x + Inches(0.3), col_y + Inches(1.15),
         col_w - Inches(0.6), Inches(0.9), fill=WHITE)
add_text(s, "d₁ᵐⁱⁿ   <   d_rᵐᵃʲ   ?",
         col_r_x + Inches(0.3), col_y + Inches(1.15),
         col_w - Inches(0.6), Inches(0.9),
         size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
add_text(s, "Majority rank shifted to r = N_maj / N_min",
         col_r_x + Inches(0.3), col_y + Inches(2.2),
         col_w - Inches(0.6), Inches(0.3),
         size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
add_rect(s, col_r_x + Inches(0.3), col_y + Inches(2.65),
         col_w - Inches(0.6), Inches(0.7), fill=PALE)
add_text(s, "→  unbiased under Poisson sampling",
         col_r_x + Inches(0.3), col_y + Inches(2.65),
         col_w - Inches(0.6), Inches(0.7),
         size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)

add_text(s, "FAMILY OF VARIANTS",
         Inches(0.6), Inches(5.6), Inches(12), Inches(0.3),
         size=11, bold=True, color=NAVY, font=FONT_HEAD)
variants = [
    ("KNNFairRank",          "core:  kₑff = r"),
    ("KNNFairRankCV",        "α tuned in  kₑff = rᵅ"),
    ("KNNFairRankJointCV",   "joint CV  (α, n_votes)"),
    ("KNNFairRankEnsemble",  "average over full grid"),
    ("KNNFairRankJackknife", "LOO variance reduction"),
]
v_w = Inches(2.45); v_h = Inches(1.05); v_y = Inches(5.95)
for i, (name, desc) in enumerate(variants):
    vx = Inches(0.6 + i * 2.5)
    add_rect(s, vx, v_y, v_w, v_h, fill=LIGHT)
    add_rect(s, vx, v_y, v_w, Inches(0.04), fill=NAVY)
    add_text(s, name, vx + Inches(0.15), v_y + Inches(0.13),
             v_w, Inches(0.35),
             size=11, bold=True, color=NAVY, font=FONT_HEAD)
    add_text(s, desc, vx + Inches(0.15), v_y + Inches(0.48),
             v_w - Inches(0.3), Inches(0.5),
             size=10, color=DARK)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — Experimental setup  (compacted: 5 essential rows)
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Setup em cinco linhas. 49 datasets binários, dos quais 40 entram na "
    "análise principal. 50 splits por par (algoritmo, dataset). Inner "
    "3-fold para tuning de α e n_votes. Cinco métricas — threshold-based "
    "e ranking-based. E uma battery estatística completa de Demšar.\n\n"
    "Tudo seedado com 42 em settings.yaml; resultados lidos de "
    "benchmark_5rep.csv, totalmente reproduzível.\n\n"
    "[≈35s. Transição: antes dos resultados, mostramos como escolhemos "
    "estes 10 algoritmos.]"
)
slide_header(s, "EXPERIMENTAL DESIGN", "Benchmark configuration", 7)

rows = [
    ("Datasets",          "49 binary CI datasets   ·   40 non-degenerate (main)   ·   9 degenerate (annex)"),
    ("Cross-validation",  "Outer: 5 reps × 10-fold StratifiedKFold = 50 splits   ·   Inner: 3-fold for (α, n_votes)"),
    ("Metrics",           "Threshold  —  G-Mean · MCC · F1     Ranking  —  ROC-AUC · PR-AUC"),
    ("Hyperparameters",   "α-grid {0.25, 0.50, 0.75, 1.00}   ·   n_votes grid {1, 2, 3, 5, 7, 10}"),
    ("Stats & repro",     "Friedman → Wilcoxon-Holm → Cohen's d   ·   seed = 42   ·   benchmark_5rep.csv (24,500 rows)"),
]
table_x, table_y = Inches(0.55), Inches(2.1)
table_w = Inches(12.2); row_h = Inches(0.85); gap = Inches(0.12)
left_w = Inches(2.8)
for i, (k, v) in enumerate(rows):
    y = table_y + i * (row_h + gap)
    add_rect(s, table_x, y, table_w, row_h, fill=LIGHT)
    add_rect(s, table_x, y, Inches(0.08), row_h, fill=NAVY)
    add_text(s, k, table_x + Inches(0.25), y, left_w, row_h,
             size=14, bold=True, color=NAVY,
             valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
    add_text(s, v, table_x + left_w + Inches(0.25), y,
             table_w - left_w - Inches(0.3), row_h,
             size=13, color=DARK, valign=MSO_ANCHOR.MIDDLE)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — Candidate screening 27 → 10
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Antes da benchmark completa fizemos screening de 1-rep com 27 "
    "algoritmos. As barras azuis são os 10 retidos; as cinzentas os "
    "17 descartados.\n\n"
    "Cada um dos 17 cai numa de quatro famílias conceptuais: paradigma "
    "errado, redundância com α-CV, density estimation ruidosa, ou "
    "dominância estrita. Detalhe no slide de anexo.\n\n"
    "[≈30s. Transição: agora aos resultados.]"
)
slide_header(s, "CANDIDATE SCREENING", "From 27 algorithms to 10", 8)

add_image(s, FIG_DIR / "screening_27_algorithms.png",
          Inches(0.55), Inches(1.75), w=Inches(8.0))

cap_x, cap_w = Inches(8.85), Inches(4.1)
add_rect(s, cap_x, Inches(1.85), cap_w, Inches(4.6), fill=LIGHT)
add_rect(s, cap_x, Inches(1.85), Inches(0.08), Inches(4.6), fill=NAVY)
add_text(s, "ELIMINATION GROUPS",
         cap_x + Inches(0.25), Inches(2.0), cap_w, Inches(0.3),
         size=11, bold=True, color=NAVY, font=FONT_HEAD)

groups = [
    ("G1 · KNNAdaptive*",         "4  —  wrong paradigm\n(adapt metric, not vote)"),
    ("G2 · Magnitude / LocalOdds", "5  —  redundant with\nα-CV"),
    ("G3 · Density / Bayesian",    "3  —  noisy estimation\nin high dim"),
    ("G4 · Topo / Ensemble",       "5  —  dominated by\nTopoJointBootstrap"),
]
for i, (k, v) in enumerate(groups):
    y = Inches(2.35 + i * 1.0)
    add_text(s, k, cap_x + Inches(0.25), y,
             cap_w - Inches(0.4), Inches(0.3),
             size=12, bold=True, color=DARK, font=FONT_HEAD)
    add_text(s, v, cap_x + Inches(0.25), y + Inches(0.3),
             cap_w - Inches(0.4), Inches(0.65),
             size=11, color=GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — Headline results: 3-metric compact table
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Tabela master compacta: G-Mean para decisão hard, ROC-AUC para "
    "ranking, MCC como referência clássica.\n\n"
    "A célula destacada a azul escuro no canto superior é o KNNFairRankJointCV "
    "em G-Mean — média 0.80 contra 0.67 do KNNOptK. A célula destacada na "
    "segunda linha é o KNNFairRankEnsemble em ROC-AUC — 0.878 contra 0.84 do "
    "SMOTE+KNN.\n\n"
    "Não há um único vencedor — voltaremos a este trade-off nas conclusões.\n\n"
    "[≈40s. Transição: vamos olhar para o G-Mean isoladamente.]"
)
slide_header(s, "RESULTS", "Headline performance — 40 datasets × 50 splits", 9)

# 3-metric table, ranked by G-Mean. Highlight champions per metric.
ALGO_ROWS_3 = [
    # (name, g_mean, roc_auc, mcc, is_champion_gmean, is_champion_roc, is_champion_mcc)
    ("KNNFairRankJointCV",            "0.799", "0.871", "0.531", True,  False, True),
    ("KNNFairRankEnsemble",           "0.781", "0.878", "0.516", False, True,  False),
    ("KNNFairRankCV",                 "0.788", "0.870", "0.527", False, False, False),
    ("KNNFairRankOptVotes",           "0.785", "0.866", "0.522", False, False, False),
    ("KNNFairRank",                   "0.775", "0.860", "0.510", False, False, False),
    ("KNNFairRankJackknife",          "0.772", "0.872", "0.510", False, False, False),
    ("KNNFairRankTopoJointBootstrap", "0.768", "0.860", "0.503", False, False, False),
    ("SMOTE+KNN",                     "0.755", "0.844", "0.482", False, False, False),
    ("KNNWeighted",                   "0.689", "0.812", "0.420", False, False, False),
    ("KNNOptK",                       "0.673", "0.840", "0.440", False, False, False),
]
metrics = [("G-Mean", "threshold"), ("ROC-AUC", "ranking"), ("MCC", "threshold")]
t_x, t_y = Inches(1.5), Inches(1.7)
col_alg_w = Inches(4.6); col_m_w = Inches(1.9); row_h = Inches(0.45)

# Super-header (metric type)
sub_x = t_x + col_alg_w
add_rect(s, t_x, t_y, col_alg_w, row_h, fill=NAVY)
add_text(s, "Algorithm  (sorted by G-Mean)", t_x + Inches(0.2), t_y,
         col_alg_w, row_h,
         size=11, bold=True, color=WHITE,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
for j, (m, kind) in enumerate(metrics):
    x = sub_x + j * col_m_w
    add_rect(s, x, t_y, col_m_w, row_h, fill=NAVY)
    add_text(s, m, x, t_y, col_m_w, row_h,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)

# Body rows
for i, row in enumerate(ALGO_ROWS_3):
    name, g, r, m, cg, cr, cm = row
    y = t_y + (i + 1) * row_h
    bg = WHITE if i % 2 == 0 else LIGHT
    add_rect(s, t_x, y, col_alg_w, row_h, fill=bg)
    add_text(s, name, t_x + Inches(0.2), y, col_alg_w, row_h,
             size=11, color=DARK, bold=(cg or cr or cm),
             valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
    for j, (val, is_champ) in enumerate([(g, cg), (r, cr), (m, cm)]):
        x = sub_x + j * col_m_w
        cell_bg = NAVY if is_champ else bg
        cell_fg = WHITE if is_champ else DARK
        add_rect(s, x, y, col_m_w, row_h, fill=cell_bg)
        add_text(s, val, x, y, col_m_w, row_h,
                 size=11, color=cell_fg, bold=is_champ,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
                 font=FONT_HEAD)

# Caption immediately under the title (replaces bottom callout that overlapped table)
add_text(s, "Champion cells highlighted in navy   ·   JointCV  →  threshold metrics   ·   Ensemble  →  ranking metric",
         Inches(0.55), Inches(1.45), Inches(12.2), Inches(0.25),
         size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 — G-Mean champion gap
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "G-Mean isolado. Em cima o KNNFairRankJointCV com 0.80; em baixo o "
    "KNNOptK com 0.67. Diferença de 0.13 — grande para esta métrica.\n\n"
    "À direita o efeito quantificado: Cohen's d pareado de +0.88 — large "
    "effect pelo critério de Cohen. Wilcoxon-Holm corrigido a 9 "
    "comparações dá p < 10⁻⁴, e o bootstrap CI da diferença não cruza zero.\n\n"
    "[≈40s. Transição: e ao longo do espectro de imbalance?]"
)
slide_header(s, "RESULTS · G-Mean", "Champion gap vs principled baseline", 10)

chart_x = Inches(0.55); chart_y = Inches(1.8); chart_w = Inches(7.8); chart_h = Inches(4.7)
add_rect(s, chart_x, chart_y, chart_w, chart_h, fill=LIGHT)
algo_gmeans = [
    ("KNNFairRankJointCV",   0.799, NAVY),
    ("KNNFairRankCV",        0.788, BLUE),
    ("KNNFairRankOptVotes",  0.785, BLUE),
    ("KNNFairRankEnsemble",  0.781, BLUE),
    ("KNNFairRank",          0.775, BLUE),
    ("KNNFairRankJackknife", 0.772, BLUE),
    ("KNNFairRankTopoJB",    0.768, BLUE),
    ("SMOTE+KNN",            0.755, GRAY),
    ("KNNWeighted",          0.689, GRAY),
    ("KNNOptK",              0.673, GRAY),
]
bar_y0 = chart_y + Inches(0.35); bar_h = Inches(0.30); bar_gap = Inches(0.10)
label_w = Inches(2.1)
for i, (name, val, col) in enumerate(algo_gmeans):
    y = bar_y0 + i * (bar_h + bar_gap)
    add_text(s, name, chart_x + Inches(0.2), y - Inches(0.04),
             label_w, bar_h + Inches(0.08),
             size=10, color=DARK, valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
    bw = Inches(val * 4.5)
    add_rect(s, chart_x + label_w + Inches(0.2), y, bw, bar_h, fill=col)
    add_text(s, f"{val:.3f}",
             chart_x + label_w + Inches(0.25) + bw, y - Inches(0.04),
             Inches(0.8), bar_h + Inches(0.08),
             size=10, bold=True, color=DARK, valign=MSO_ANCHOR.MIDDLE,
             font=FONT_HEAD)
add_text(s, "Mean G-Mean across 40 datasets",
         chart_x, chart_y + Inches(4.25), chart_w, Inches(0.3),
         size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

right_x = Inches(8.7); right_w = Inches(4.2)
add_rect(s, right_x, Inches(1.9), right_w, Inches(2.0), fill=NAVY)
add_text(s, "Cohen's d", right_x, Inches(2.0), right_w, Inches(0.4),
         size=12, color=PALE, align=PP_ALIGN.CENTER, font=FONT_HEAD, bold=True)
add_text(s, "+0.88", right_x, Inches(2.4), right_w, Inches(1.2),
         size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
add_text(s, "LARGE effect", right_x, Inches(3.55), right_w, Inches(0.3),
         size=12, color=PALE, italic=True, align=PP_ALIGN.CENTER)
y0 = Inches(4.1)
add_text(s, "JointCV  vs  KNNOptK", right_x, y0, right_w, Inches(0.3),
         size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font=FONT_HEAD)
stats = [
    ("Δ mean G-Mean", "+0.126"),
    ("Wilcoxon raw p", "< 10⁻⁴"),
    ("Holm-corrected p", "< 10⁻⁴"),
    ("Bootstrap 95% CI", "[+0.10, +0.15]"),
]
for i, (k, v) in enumerate(stats):
    y = Inches(4.5 + i * 0.45)
    add_text(s, k, right_x, y, right_w / 2 + Inches(0.6), Inches(0.4),
             size=11, color=GRAY, valign=MSO_ANCHOR.MIDDLE)
    add_text(s, v, right_x + right_w / 2 + Inches(0.6), y,
             right_w / 2 - Inches(0.6), Inches(0.4),
             size=11, bold=True, color=DARK, align=PP_ALIGN.RIGHT,
             valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 — Universal dominance across IR
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Pergunta natural: 'a vantagem do FairRank vem só dos datasets extremos, "
    "ou aparece em todo o espectro?'.\n\n"
    "Dividimos os 40 datasets em quartis por IR. A linha tracejada cinzenta "
    "é o KNNOptK. As linhas a cores são os FairRank. Em TODOS os quartis "
    "— incluindo Q1 com IR baixo — o FairRank fica acima do baseline.\n\n"
    "Kendall τ entre wins e log-IR: -0.08, p = 0.54. NÃO há tendência "
    "monótona, mas o win rate é 85%. Universal dominance — não condicional "
    "a IR.\n\n"
    "[≈45s. Transição: validação estatística formal.]"
)
slide_header(s, "RESULTS · IR quartile analysis", "Advantage is universal, not IR-conditional", 11)

add_image(s, FIG_DIR / "bias_by_ir.png", Inches(0.55), Inches(1.7), w=Inches(8.5))

cap_x = Inches(9.4); cap_w = Inches(3.6)
add_rect(s, cap_x, Inches(1.85), cap_w, Inches(1.85), fill=NAVY)
add_text(s, "JointCV win rate", cap_x, Inches(1.95), cap_w, Inches(0.3),
         size=11, color=PALE, align=PP_ALIGN.CENTER, font=FONT_HEAD, bold=True)
add_text(s, "85%", cap_x, Inches(2.3), cap_w, Inches(1.0),
         size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
add_text(s, "of 40 datasets", cap_x, Inches(3.35), cap_w, Inches(0.3),
         size=11, color=PALE, italic=True, align=PP_ALIGN.CENTER)

add_rect(s, cap_x, Inches(3.95), cap_w, Inches(1.45), fill=LIGHT)
add_rect(s, cap_x, Inches(3.95), Inches(0.08), Inches(1.45), fill=NAVY)
add_text(s, "Kendall trend test",
         cap_x + Inches(0.25), Inches(4.1),
         cap_w, Inches(0.3),
         size=11, bold=True, color=NAVY, font=FONT_HEAD)
add_text(s, "τ = −0.08    p = 0.54",
         cap_x + Inches(0.25), Inches(4.4), cap_w, Inches(0.4),
         size=15, bold=True, color=DARK, font=FONT_HEAD)
add_text(s, "Flat trend across IR axis",
         cap_x + Inches(0.25), Inches(4.85), cap_w, Inches(0.5),
         size=11, italic=True, color=GRAY)

add_rect(s, cap_x, Inches(5.6), cap_w, Inches(1.0), fill=PALE)
add_text(s, "Universal, not IR-conditional.",
         cap_x + Inches(0.2), Inches(5.7), cap_w - Inches(0.3), Inches(0.85),
         size=13, bold=True, color=NAVY, font=FONT_HEAD,
         valign=MSO_ANCHOR.MIDDLE)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12 — Statistical validation (Demšar)  — simplified
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Validação estatística completa, protocolo de Demšar 2006.\n\n"
    "Primeiro Friedman global: rejeitamos a nula 'todos iguais' com "
    "p < 10⁻¹² em todas as 5 métricas. Depois post-hoc Wilcoxon "
    "signed-rank pareado com correcção Holm, contra dois baselines.\n\n"
    "Linha 1: JointCV vs KNNOptK em G-Mean — p_holm < 10⁻⁴, d = +0.88, "
    "large effect. Linha 2: Ensemble vs SMOTE+KNN em ROC-AUC — "
    "p_holm < 10⁻⁴, d = +0.56, medium-large.\n\n"
    "[≈40s. Transição: a teoria que motiva a correcção — testámo-la "
    "empiricamente?]"
)
slide_header(s, "RESULTS · Statistical validation", "Demšar 2006 — rejected, large, robust", 12)

# Friedman summary panel
y_top = Inches(1.85); y_top_h = Inches(1.4)
add_rect(s, Inches(0.55), y_top, Inches(12.25), y_top_h, fill=LIGHT)
add_rect(s, Inches(0.55), y_top, Inches(0.12), y_top_h, fill=NAVY)
add_text(s, "① FRIEDMAN  ·  global test",
         Inches(0.75), y_top + Inches(0.15), Inches(4), Inches(0.3),
         size=11, bold=True, color=NAVY, font=FONT_HEAD)
add_text(s, "Rejected on all 5 metrics  —  p < 10⁻¹²",
         Inches(0.75), y_top + Inches(0.45), Inches(7), Inches(0.5),
         size=20, bold=True, color=DARK, font=FONT_HEAD)
add_text(s, "Algorithms differ; warrants post-hoc analysis.",
         Inches(0.75), y_top + Inches(0.95), Inches(11), Inches(0.4),
         size=12, italic=True, color=GRAY)

# Post-hoc table
y_post = Inches(3.55); rh = Inches(0.6); table_w_pst = Inches(12.25)
# header
add_rect(s, Inches(0.55), y_post, table_w_pst, rh, fill=NAVY)
add_text(s, "② POST-HOC  ·  Wilcoxon signed-rank  +  Holm correction  +  Cohen's d",
         Inches(0.55), y_post, table_w_pst, rh,
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)

# rows
posthoc = [
    ("JointCV",  "vs  KNNOptK",    "G-Mean",   "< 10⁻⁴", "+0.88", "LARGE",  GREEN),
    ("Ensemble", "vs  SMOTE+KNN",  "ROC-AUC",  "< 10⁻⁴", "+0.56", "MEDIUM", AMBER),
]
col_alg = Inches(2.0); col_vs = Inches(2.2); col_met = Inches(1.8)
col_p   = Inches(2.0); col_d   = Inches(1.8); col_eff = Inches(2.45)
xs = [Inches(0.55), Inches(0.55)+col_alg, Inches(0.55)+col_alg+col_vs,
      Inches(0.55)+col_alg+col_vs+col_met,
      Inches(0.55)+col_alg+col_vs+col_met+col_p,
      Inches(0.55)+col_alg+col_vs+col_met+col_p+col_d]
ws = [col_alg, col_vs, col_met, col_p, col_d, col_eff]
labels_hdr = ["Algorithm", "Comparison", "Metric", "p_holm", "Cohen's d", "Effect"]
y_hdr = y_post + rh
add_rect(s, Inches(0.55), y_hdr, table_w_pst, rh, fill=LIGHT)
for j, lbl in enumerate(labels_hdr):
    add_text(s, lbl, xs[j], y_hdr, ws[j], rh,
             size=10, bold=True, color=GRAY, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)

for i, (alg, vs_, met, p, d, eff, ecol) in enumerate(posthoc):
    y = y_hdr + (i + 1) * rh
    add_rect(s, Inches(0.55), y, table_w_pst, rh, fill=WHITE)
    add_text(s, alg, xs[0], y, ws[0], rh,
             size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
    add_text(s, vs_, xs[1], y, ws[1], rh,
             size=12, color=DARK, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
    add_text(s, met, xs[2], y, ws[2], rh,
             size=12, color=DARK, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
    add_text(s, p, xs[3], y, ws[3], rh,
             size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
    add_text(s, d, xs[4], y, ws[4], rh,
             size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
    add_rect(s, xs[5] + Inches(0.4), y + Inches(0.12),
             ws[5] - Inches(0.8), rh - Inches(0.24), fill=ecol)
    add_text(s, eff, xs[5], y, ws[5], rh,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 13 — Empirical Poisson validation
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Testámos empiricamente a assunção Poisson que motiva o algoritmo. "
    "Para cada dataset, comparámos a razão empírica d_1^min / d_1^maj "
    "com o teórico r^(1/d).\n\n"
    "No scatter log-log à direita, 37 dos 40 datasets caem ACIMA da "
    "linha de acordo perfeito. As minorities estão MAIS clusterizadas "
    "que um Poisson homogéneo previa.\n\n"
    "Isto não invalida a derivação — motiva o expoente α < 1 que o "
    "KNNFairRankCV introduz: se a violação é para mais clustering, "
    "amorteça-se a correcção.\n\n"
    "[≈50s. Transição: e isto prediz a vantagem do FairRank?]"
)
slide_header(s, "RESULTS · Theory check", "Poisson holds — but minorities cluster more than predicted", 13)

add_image(s, FIG_DIR / "poisson_empirical_validation.png",
          Inches(0.55), Inches(1.7), w=Inches(9.2))

cap_x = Inches(10.0); cap_w = Inches(3.0)
add_rect(s, cap_x, Inches(1.85), cap_w, Inches(1.6), fill=NAVY)
add_text(s, "Empirical > Theory", cap_x, Inches(1.95), cap_w, Inches(0.3),
         size=10, color=PALE, align=PP_ALIGN.CENTER, font=FONT_HEAD, bold=True)
add_text(s, "37 / 40", cap_x, Inches(2.25), cap_w, Inches(0.95),
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
add_text(s, "datasets clustered", cap_x, Inches(3.1), cap_w, Inches(0.3),
         size=10, color=PALE, italic=True, align=PP_ALIGN.CENTER)

add_rect(s, cap_x, Inches(3.65), cap_w, Inches(1.6), fill=LIGHT)
add_rect(s, cap_x, Inches(3.65), Inches(0.08), Inches(1.6), fill=NAVY)
add_text(s, "Implication",
         cap_x + Inches(0.25), Inches(3.8), cap_w, Inches(0.3),
         size=11, bold=True, color=NAVY, font=FONT_HEAD)
add_text(s, "Motivates α < 1\nin KNNFairRankCV —\ndampen the correction.",
         cap_x + Inches(0.25), Inches(4.15),
         cap_w - Inches(0.4), Inches(1.0),
         size=12, color=DARK)

add_rect(s, cap_x, Inches(5.45), cap_w, Inches(1.2), fill=PALE)
add_text(s, "Median log₂ deviation",
         cap_x + Inches(0.2), Inches(5.55), cap_w - Inches(0.3), Inches(0.3),
         size=10, color=GRAY, align=PP_ALIGN.CENTER, italic=True)
add_text(s, "≈ +1.0",
         cap_x + Inches(0.2), Inches(5.85), cap_w - Inches(0.3), Inches(0.7),
         size=26, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14 — Theory ↔ data link (the killer slide)
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Esta é a slide mais importante. A teoria faz uma predição falsificável: "
    "o FairRank deve ajudar mais onde o KNN é mais enviesado — minorities "
    "na fronteira de decisão.\n\n"
    "Computámos duas meta-features estruturais: borderline_fraction (fracção "
    "de minorities cujo nearest neighbour é majority) e minority_nbr_purity "
    "(pureza dos 10-NN dos minorities). Correlacionámos com Δ G-Mean do "
    "JointCV via Spearman.\n\n"
    "Resultados à direita: ρ = +0.77 e ρ = -0.76, ambos p < 10⁻⁴. As DUAS "
    "direcções coincidem com o que a teoria prediz. Direct theory-data link.\n\n"
    "[≈55s. Transição: discussão e conclusões.]"
)
slide_header(s, "THEORY ↔ DATA LINK", "Theory predicts where to win — empirically confirmed", 14)

add_image(s, FIG_DIR / "knn_meta_vs_advantage.png",
          Inches(0.55), Inches(1.7), w=Inches(8.5))

cap_x = Inches(9.4); cap_w = Inches(3.6)
add_text(s, "Theory predicts",
         cap_x, Inches(1.8), cap_w, Inches(0.3),
         size=11, bold=True, color=NAVY, font=FONT_HEAD)
add_rect(s, cap_x, Inches(2.15), cap_w, Inches(2.05), fill=LIGHT)
add_rect(s, cap_x, Inches(2.15), Inches(0.08), Inches(2.05), fill=GREEN)
add_text(s, "borderline_fraction  ↑",
         cap_x + Inches(0.25), Inches(2.3),
         cap_w - Inches(0.4), Inches(0.3),
         size=11, bold=True, color=GREEN, font=FONT_HEAD)
add_text(s, "more minorities at the\nboundary → larger Δ",
         cap_x + Inches(0.25), Inches(2.6),
         cap_w - Inches(0.4), Inches(0.7),
         size=11, color=DARK, italic=True)
add_text(s, "ρ = +0.77",
         cap_x + Inches(0.25), Inches(3.3),
         cap_w - Inches(0.4), Inches(0.45),
         size=20, bold=True, color=NAVY, font=FONT_HEAD)
add_text(s, "p < 10⁻⁴   (n = 38)",
         cap_x + Inches(0.25), Inches(3.7),
         cap_w - Inches(0.4), Inches(0.4),
         size=11, color=GRAY)

add_rect(s, cap_x, Inches(4.35), cap_w, Inches(2.05), fill=LIGHT)
add_rect(s, cap_x, Inches(4.35), Inches(0.08), Inches(2.05), fill=GREEN)
add_text(s, "minority_nbr_purity  ↑",
         cap_x + Inches(0.25), Inches(4.5),
         cap_w - Inches(0.4), Inches(0.3),
         size=11, bold=True, color=GREEN, font=FONT_HEAD)
add_text(s, "purer minority clusters\n→ smaller Δ",
         cap_x + Inches(0.25), Inches(4.8),
         cap_w - Inches(0.4), Inches(0.7),
         size=11, color=DARK, italic=True)
add_text(s, "ρ = −0.76",
         cap_x + Inches(0.25), Inches(5.5),
         cap_w - Inches(0.4), Inches(0.45),
         size=20, bold=True, color=NAVY, font=FONT_HEAD)
add_text(s, "p < 10⁻⁴   (n = 38)",
         cap_x + Inches(0.25), Inches(5.9),
         cap_w - Inches(0.4), Inches(0.4),
         size=11, color=GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 15 — Trade-off discovered
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Achado interessante antes de concluir. Não há um único algoritmo "
    "vencedor — o trade-off cinde-se por tipo de métrica.\n\n"
    "À esquerda, threshold metrics: JointCV vence — escolhe um par óptimo "
    "(α, n_votes) via inner CV, dando decisões crisp. À direita, ranking "
    "metrics: Ensemble vence — faz média sobre o α-grid, dando "
    "probabilidades suaves que ranqueiam melhor.\n\n"
    "O trade-off é matematicamente coerente: tarefas diferentes (0-1 vs "
    "AUC) têm losses diferentes, e é esperado que predictores distintos "
    "sejam óptimos para cada uma.\n\n"
    "[≈40s. Transição: conclusões.]"
)
slide_header(s, "DISCUSSION", "Trade-off — mathematically coherent split", 15)

col_y = Inches(1.85); col_h = Inches(4.7)
col_w = Inches(5.9); col_l_x = Inches(0.55); col_r_x = Inches(6.85)

add_rect(s, col_l_x, col_y, col_w, col_h, fill=LIGHT)
add_rect(s, col_l_x, col_y, col_w, Inches(0.7), fill=NAVY)
add_text(s, "THRESHOLD METRICS", col_l_x, col_y, col_w, Inches(0.4),
         size=11, bold=True, color=PALE, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
add_text(s, "G-Mean   ·   MCC",
         col_l_x, col_y + Inches(0.32), col_w, Inches(0.38),
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
add_text(s, "Winner", col_l_x + Inches(0.4), col_y + Inches(0.9),
         col_w - Inches(0.6), Inches(0.3),
         size=11, color=GRAY, font=FONT_HEAD)
add_text(s, "KNNFairRankJointCV",
         col_l_x + Inches(0.4), col_y + Inches(1.2),
         col_w - Inches(0.6), Inches(0.6),
         size=24, bold=True, color=NAVY, font=FONT_HEAD)
add_text(s, "Why", col_l_x + Inches(0.4), col_y + Inches(2.0),
         col_w - Inches(0.6), Inches(0.3),
         size=11, color=GRAY, font=FONT_HEAD)
add_text(s, "Inner CV picks a single optimal (α, n_votes) pair.",
         col_l_x + Inches(0.4), col_y + Inches(2.3),
         col_w - Inches(0.6), Inches(0.6),
         size=13, color=DARK)
add_text(s, "Crisp probabilities → sharp 0/1 decisions.",
         col_l_x + Inches(0.4), col_y + Inches(2.85),
         col_w - Inches(0.6), Inches(0.6),
         size=13, color=DARK, italic=True)
add_rect(s, col_l_x + Inches(0.4), col_y + Inches(3.6),
         col_w - Inches(0.8), Inches(0.85), fill=PALE)
add_text(s, "G-Mean 0.799    ·    MCC 0.531",
         col_l_x + Inches(0.4), col_y + Inches(3.6),
         col_w - Inches(0.8), Inches(0.85),
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)

add_rect(s, col_r_x, col_y, col_w, col_h, fill=LIGHT)
add_rect(s, col_r_x, col_y, col_w, Inches(0.7), fill=NAVY)
add_text(s, "RANKING METRICS", col_r_x, col_y, col_w, Inches(0.4),
         size=11, bold=True, color=PALE, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
add_text(s, "ROC-AUC   ·   PR-AUC",
         col_r_x, col_y + Inches(0.32), col_w, Inches(0.38),
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
add_text(s, "Winner", col_r_x + Inches(0.4), col_y + Inches(0.9),
         col_w - Inches(0.6), Inches(0.3),
         size=11, color=GRAY, font=FONT_HEAD)
add_text(s, "KNNFairRankEnsemble",
         col_r_x + Inches(0.4), col_y + Inches(1.2),
         col_w - Inches(0.6), Inches(0.6),
         size=24, bold=True, color=NAVY, font=FONT_HEAD)
add_text(s, "Why", col_r_x + Inches(0.4), col_y + Inches(2.0),
         col_w - Inches(0.6), Inches(0.3),
         size=11, color=GRAY, font=FONT_HEAD)
add_text(s, "Averages vote fractions across the full α-grid.",
         col_r_x + Inches(0.4), col_y + Inches(2.3),
         col_w - Inches(0.6), Inches(0.6),
         size=13, color=DARK)
add_text(s, "Smoother probabilities → better ranking quality.",
         col_r_x + Inches(0.4), col_y + Inches(2.85),
         col_w - Inches(0.6), Inches(0.6),
         size=13, color=DARK, italic=True)
add_rect(s, col_r_x + Inches(0.4), col_y + Inches(3.6),
         col_w - Inches(0.8), Inches(0.85), fill=PALE)
add_text(s, "ROC-AUC 0.878    ·    PR-AUC 0.649",
         col_r_x + Inches(0.4), col_y + Inches(3.6),
         col_w - Inches(0.8), Inches(0.85),
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)

add_text(s, "Mathematically coherent: different losses → different optimal predictors.",
         Inches(0.55), Inches(6.75), Inches(12.2), Inches(0.4),
         size=12, bold=True, color=NAVY, italic=True,
         align=PP_ALIGN.CENTER, font=FONT_HEAD)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 16 — Conclusions & Future Work
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Para fechar.\n\n"
    "Contribuições: uma correcção ao KNN derivada de princípios — não "
    "heurística; validação rigorosa em 49 datasets sob Demšar; e um "
    "direct theory-data link com ρ = +0.77.\n\n"
    "Limitações: a assunção Poisson é violada em ~95% dos datasets, "
    "tratada via α < 1 mas extreme clustering ainda é difícil; o "
    "k_maj_cap satura em 1 dos 40 datasets; e o multi-vote scheme "
    "ainda não tem justificação formal completa.\n\n"
    "Future work: extensão Neyman-Scott para clustered minorities, "
    "generalização multi-class, e comparação directa com DANN.\n\n"
    "Obrigado. Estamos disponíveis para questões.\n\n"
    "[≈60s. FIM da apresentação principal — total ~10:30 minutos, "
    "deixa ~1:30 de buffer para Q&A.]"
)
slide_header(s, "CONCLUSIONS", "Contributions  ·  Limitations  ·  Future Work", 16)

col_y = Inches(1.85); col_h = Inches(4.8)
col_w = Inches(4.1); col_gap = Inches(0.15)
c1_x = Inches(0.55); c2_x = c1_x + col_w + col_gap; c3_x = c2_x + col_w + col_gap

def col_block(x, accent, kicker, items):
    add_rect(s, x, col_y, col_w, col_h, fill=LIGHT)
    add_rect(s, x, col_y, col_w, Inches(0.6), fill=accent)
    add_text(s, kicker, x, col_y + Inches(0.12),
             col_w, Inches(0.4),
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
    for i, item in enumerate(items):
        y = col_y + Inches(0.85 + i * 1.25)
        bullet = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    x + Inches(0.3), y + Inches(0.06),
                                    Inches(0.22), Inches(0.22))
        bullet.fill.solid(); bullet.fill.fore_color.rgb = accent
        bullet.line.fill.background(); bullet.shadow.inherit = False
        add_text(s, str(i + 1),
                 x + Inches(0.3), y + Inches(0.06),
                 Inches(0.22), Inches(0.22),
                 size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
        add_text(s, item,
                 x + Inches(0.65), y,
                 col_w - Inches(0.8), Inches(1.1),
                 size=12, color=DARK)

col_block(c1_x, GREEN, "CONTRIBUTIONS", [
    "Principle-derived correction:\nkₑff = r from Poisson order statistics",
    "Rigorous validation:\n49 ds × 50 splits, full Demšar 2006",
    "Direct theory ↔ data link:\nρ = +0.77 with borderline_fraction",
])
col_block(c2_x, AMBER, "LIMITATIONS", [
    "Poisson assumption violated\non ~95% of datasets (handled by α-CV)",
    "k_maj_cap saturates on 1/40\ndatasets at extreme IR",
    "Multi-vote scheme not yet\nformally justified",
])
col_block(c3_x, BLUE, "FUTURE WORK", [
    "Neyman-Scott extension\nfor clustered-minority theory",
    "Multi-class generalisation\n+ regression analogue",
    "Direct comparison\nwith DANN  ·  threshold tuning",
])


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 17 — Annex divider
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "[Slide divisor de anexos. Mostrar apenas se houver perguntas técnicas "
    "no Q&A.]"
)
add_rect(s, Inches(0), Inches(0), SW, SH, fill=NAVY)
add_text(s, "ANNEXES",
         Inches(0), Inches(3.0), SW, Inches(0.6),
         size=14, bold=True, color=PALE, align=PP_ALIGN.CENTER, font=FONT_HEAD)
add_text(s, "Supporting material",
         Inches(0), Inches(3.55), SW, Inches(1.2),
         size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_HEAD)
add_text(s, "For Q&A and technical discussion",
         Inches(0), Inches(4.7), SW, Inches(0.5),
         size=16, color=PALE, italic=True, align=PP_ALIGN.CENTER)
add_rect(s, Inches(5.665), Inches(5.5), Inches(2.0), Inches(0.04), fill=BLUE)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 18 — ANNEX: rank convergence (moved from main flow)
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Rank-convergence diagnostic: para cada n ∈ {5, 10, ..., 50} splits "
    "acumulados recalculámos os ranks médios. Para os campeões (JointCV, "
    "Ensemble), |Δrank| entre n=40 e n=50 é ≤ 0.075 — ordem de grandeza "
    "abaixo dos rank gaps típicos. 50 splits são suficientes.\n\n"
    "Q&A típico: 'porque 50 e não 100?' — porque a partir de 35 splits "
    "os ranks dos campeões já estão dentro de 0.1 da posição final."
)
slide_header(s, "ANNEX · Robustness", "Are 50 splits enough?  —  rank-convergence", 18)

add_image(s, FIG_DIR / "rank_convergence.png", Inches(0.55), Inches(1.7),
          w=Inches(8.5))

cap_x = Inches(9.4); cap_w = Inches(3.6)
add_rect(s, cap_x, Inches(1.85), cap_w, Inches(1.65), fill=NAVY)
add_text(s, "max |Δrank|  n=40→50", cap_x, Inches(1.95), cap_w, Inches(0.3),
         size=10, color=PALE, align=PP_ALIGN.CENTER, font=FONT_HEAD, bold=True)
add_text(s, "≤ 0.075", cap_x, Inches(2.25), cap_w, Inches(0.95),
         size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
add_text(s, "for both champions (JointCV, Ensemble)",
         cap_x, Inches(3.1), cap_w, Inches(0.4),
         size=10, color=PALE, italic=True, align=PP_ALIGN.CENTER)

add_rect(s, cap_x, Inches(3.75), cap_w, Inches(1.5), fill=LIGHT)
add_rect(s, cap_x, Inches(3.75), Inches(0.08), Inches(1.5), fill=NAVY)
add_text(s, "Context", cap_x + Inches(0.25), Inches(3.9), cap_w, Inches(0.3),
         size=11, bold=True, color=NAVY, font=FONT_HEAD)
add_text(s, "Typical gap between\nadjacent ranks: 1 – 2 positions",
         cap_x + Inches(0.25), Inches(4.2),
         cap_w - Inches(0.4), Inches(1.0),
         size=12, color=DARK)
add_rect(s, cap_x, Inches(5.45), cap_w, Inches(1.2), fill=PALE)
add_text(s, "Conclusion: 50 splits suffice — results are not finite-sample artefacts.",
         cap_x + Inches(0.2), Inches(5.55), cap_w - Inches(0.3), Inches(1.0),
         size=11, bold=True, color=NAVY, font=FONT_HEAD,
         valign=MSO_ANCHOR.MIDDLE)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 19 — ANNEX: ISA PCA projection
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "ISA PCA projection — 40 datasets em 2D via PCA sobre 6 meta-features. "
    "PC1 e PC2 explicam ~70% da variância. Tamanho ∝ severidade do "
    "imbalance."
)
slide_header(s, "ANNEX · ISA", "Instance space — PCA projection", 19)
add_image(s, FIG_DIR / "isa_projection.png",
          Inches(3.0), Inches(1.7), h=Inches(5.0))


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 20 — ANNEX: ISA footprints
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Algorithm footprints — para cada métrica, marcamos no instance space "
    "qual algoritmo vence cada dataset. JointCV domina G-Mean / MCC; "
    "Ensemble domina ROC-AUC / PR-AUC."
)
slide_header(s, "ANNEX · ISA", "Algorithm footprints — where each wins", 20)
add_image(s, FIG_DIR / "isa_footprints.png",
          Inches(0.55), Inches(1.7), w=Inches(12.2))


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 21 — ANNEX: ISA stress test
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Stress test ao longo do eixo IR — wins por quartil. JointCV mantém "
    "≥20% de wins em todos os quartis e cresce em Q4."
)
slide_header(s, "ANNEX · ISA", "Stress test along IR axis", 21)
add_image(s, FIG_DIR / "isa_stress_test.png",
          Inches(0.55), Inches(1.7), w=Inches(12.2))


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 22 — ANNEX: degenerate datasets
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Nove datasets com N_min < 20 foram separados da análise principal "
    "(stratified CV produz folds com 0 minorities). Em 2 destes, SMOTE "
    "colapsa para G-Mean = 0 — interpolation produz amostras sintéticas "
    "fora da minority region. FairRank baseado em counting mantém-se "
    "definido independentemente de N_min."
)
slide_header(s, "ANNEX", "Degenerate datasets  (N_min < 20)  —  exploratory only", 22)

add_rect(s, Inches(0.55), Inches(1.8), Inches(6.0), Inches(4.7), fill=LIGHT)
add_rect(s, Inches(0.55), Inches(1.8), Inches(0.08), Inches(4.7), fill=AMBER)
add_text(s, "WHY SEPARATED",
         Inches(0.8), Inches(1.95), Inches(6), Inches(0.3),
         size=11, bold=True, color=AMBER, font=FONT_HEAD)
add_text(s,
         "Stratified 10-fold CV with N_min < 20 can produce test folds with\n"
         "zero positives  →  ROC-AUC and G-Mean undefined.\n\n"
         "350 NaN rows in benchmark_5rep.csv (1.4%) — all confined\n"
         "to these 9 datasets; stripped before the main §7 analysis.",
         Inches(0.8), Inches(2.3), Inches(5.7), Inches(2.5),
         size=13, color=DARK)
add_text(s, "9 datasets",
         Inches(0.8), Inches(4.85), Inches(5.7), Inches(0.4),
         size=14, bold=True, color=NAVY, font=FONT_HEAD)
add_text(s,
         "analcatdata_challenger, kc1-top5, ar1, ar6,\n"
         "analcatdata_lawsuit, analcatdata_neavote, analcatdata_chlamydia,\n"
         "arsenic-female-lung, arsenic-male-lung",
         Inches(0.8), Inches(5.25), Inches(5.7), Inches(1.2),
         size=10, italic=True, color=GRAY)

add_rect(s, Inches(6.95), Inches(1.8), Inches(6.0), Inches(4.7), fill=LIGHT)
add_rect(s, Inches(6.95), Inches(1.8), Inches(0.08), Inches(4.7), fill=NAVY)
add_text(s, "KEY FINDING",
         Inches(7.2), Inches(1.95), Inches(6), Inches(0.3),
         size=11, bold=True, color=NAVY, font=FONT_HEAD)
add_text(s, "SMOTE collapses on 2 cases",
         Inches(7.2), Inches(2.3), Inches(5.7), Inches(0.5),
         size=18, bold=True, color=RED, font=FONT_HEAD)
add_text(s,
         "Interpolating between 7-9 minority points produces synthetic\n"
         "samples outside the true minority region  →  classifier predicts\n"
         "majority for everything  →  G-Mean = 0.",
         Inches(7.2), Inches(2.85), Inches(5.7), Inches(2.0),
         size=13, color=DARK)
add_rect(s, Inches(7.2), Inches(4.9), Inches(5.5), Inches(1.45), fill=PALE)
add_text(s, "KNNFairRank stays well-defined",
         Inches(7.3), Inches(5.0), Inches(5.3), Inches(0.4),
         size=14, bold=True, color=NAVY, font=FONT_HEAD,
         align=PP_ALIGN.CENTER)
add_text(s, "count-based correction does not rely on N_min magnitude.",
         Inches(7.3), Inches(5.45), Inches(5.3), Inches(0.85),
         size=12, color=DARK, align=PP_ALIGN.CENTER, italic=True)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 23 — ANNEX: k_maj_cap saturation
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "k_maj_cap saturation: o KNNFairRank tem um cap de 1000 majority "
    "neighbours. Apenas 1 dos 40 datasets satura — dataset_1056_mc1 com "
    "IR=138. Para os outros 39, k_eff = r é exercido na íntegra."
)
slide_header(s, "ANNEX", "k_maj_cap saturation — operational sanity", 23)

add_rect(s, Inches(2.5), Inches(2.2), Inches(8.3), Inches(2.5), fill=NAVY)
add_text(s, "1 / 40", Inches(2.5), Inches(2.4), Inches(8.3), Inches(1.4),
         size=84, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
add_text(s, "datasets saturate the cap",
         Inches(2.5), Inches(3.9), Inches(8.3), Inches(0.4),
         size=14, color=PALE, italic=True, align=PP_ALIGN.CENTER)
add_text(s, "dataset_1056_mc1   ·   IR = 138.2   ·   requested 1400   ·   capped at 1000",
         Inches(2.5), Inches(4.3), Inches(8.3), Inches(0.4),
         size=11, color=PALE, align=PP_ALIGN.CENTER)
add_text(s, "For the other 39 datasets, k_eff = r is exercised in full.",
         Inches(2.5), Inches(5.1), Inches(8.3), Inches(0.4),
         size=12, color=DARK, align=PP_ALIGN.CENTER, italic=True)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 24 — ANNEX: Wilcoxon ties
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Transparência sobre o poder estatístico. Wilcoxon descarta diferenças "
    "exactamente zero; reportamos o N efectivo por par. N_eff entre 38-39 "
    "sobre 40 — efeito pequeno, mas declarado explicitamente."
)
slide_header(s, "ANNEX", "Wilcoxon ties — effective sample size", 24)

add_rect(s, Inches(3.5), Inches(2.5), Inches(6.3), Inches(2.0), fill=NAVY)
add_text(s, "N_eff  ∈  {38, 39}", Inches(3.5), Inches(2.65), Inches(6.3),
         Inches(0.7),
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
add_text(s, "out of 40 — across all (algorithm, metric, baseline) pairs",
         Inches(3.5), Inches(3.4), Inches(6.3), Inches(0.5),
         size=12, color=PALE, italic=True, align=PP_ALIGN.CENTER)
add_text(s, "TopoJointBootstrap accounts for most ties (2 in G-Mean / MCC)",
         Inches(3.5), Inches(3.85), Inches(6.3), Inches(0.5),
         size=11, color=PALE, align=PP_ALIGN.CENTER)
add_text(s, "Tie counts are tabulated explicitly in §8.3 of the notebook.",
         Inches(0.55), Inches(5.2), Inches(12.2), Inches(0.4),
         size=12, color=DARK, align=PP_ALIGN.CENTER, italic=True)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 25 — ANNEX: complexity
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Complexidade computacional e runtime empírico. O JointCV é dominante "
    "no orçamento (84 min/dataset com 8 cores), seguido do OptVotes. "
    "Total ~2.5 horas no Apple M1 para a benchmark completa."
)
slide_header(s, "ANNEX", "Computational complexity & runtime", 25)

algos = [
    ("KNNOptK",                "O(F · k_grid · N²)",        "Inner CV over k",  "11 min"),
    ("KNNWeighted",            "O(N²)",                      "Just KNN at k = 5", "1 min"),
    ("SMOTE+KNN",              "O(N² + N_synth · k_neigh)",  "Oversample + KNN",  "1 min"),
    ("KNNFairRank",            "O(N²)",                      "Closed-form k_eff", "1 min"),
    ("KNNFairRankCV",          "O(F · α · N²)",              "Inner CV over α",   "9 min"),
    ("KNNFairRankJointCV",     "O(F · α · n_votes · N²)",    "Joint (α, n_votes) CV", "84 min"),
    ("KNNFairRankEnsemble",    "O(α · n_votes · N²)",        "Grid voting at inference", "1 min"),
    ("KNNFairRankOptVotes",    "O(F · n_votes · N²)",        "Inner CV over n_votes", "30 min"),
    ("KNNFairRankJackknife",   "O(N · k_min · N)",           "LOO over minority", "2 min"),
    ("KNNFairRankTopoJB",      "O(N · B · PH)",              "Bootstrap + topology", "5 min"),
]
tx, ty = Inches(0.55), Inches(1.7)
cw1, cw2, cw3, cw4 = Inches(3.5), Inches(3.5), Inches(3.4), Inches(1.8)
rh_ = Inches(0.43)
add_rect(s, tx, ty, cw1 + cw2 + cw3 + cw4, rh_, fill=NAVY)
add_text(s, "Algorithm", tx + Inches(0.2), ty, cw1, rh_,
         size=11, bold=True, color=WHITE, valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
add_text(s, "Complexity", tx + cw1, ty, cw2, rh_,
         size=11, bold=True, color=WHITE, valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
add_text(s, "Mechanism", tx + cw1 + cw2, ty, cw3, rh_,
         size=11, bold=True, color=WHITE, valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
add_text(s, "Time / ds", tx + cw1 + cw2 + cw3, ty, cw4, rh_,
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
for i, row in enumerate(algos):
    y = ty + (i + 1) * rh_
    bg = WHITE if i % 2 == 0 else LIGHT
    add_rect(s, tx, y, cw1 + cw2 + cw3 + cw4, rh_, fill=bg)
    add_text(s, row[0], tx + Inches(0.2), y, cw1, rh_,
             size=11, color=DARK, bold=True, valign=MSO_ANCHOR.MIDDLE,
             font=FONT_HEAD)
    add_text(s, row[1], tx + cw1, y, cw2, rh_,
             size=10, color=DARK, valign=MSO_ANCHOR.MIDDLE)
    add_text(s, row[2], tx + cw1 + cw2, y, cw3, rh_,
             size=10, color=GRAY, valign=MSO_ANCHOR.MIDDLE, italic=True)
    add_text(s, row[3], tx + cw1 + cw2 + cw3, y, cw4, rh_,
             size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
add_text(s,
         "Total benchmark runtime  ≈  2.5 h on Apple M1 (16 GB)   ·   49 datasets × 10 algorithms × 50 splits",
         Inches(0.55), Inches(6.8), Inches(12.2), Inches(0.3),
         size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 26 — ANNEX: candidate screening detail
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Detalhe dos 4 grupos eliminados no screening. Cada grupo tem "
    "motivação distinta e nenhum cobre um mecanismo de correcção que não "
    "esteja já presente nos 10 retidos."
)
slide_header(s, "ANNEX · Candidate screening", "Elimination rationale — 4 conceptual groups", 26)

groups_d = [
    ("G1", "KNNAdaptive*",          RED,
     "Adapt the metric (entropy, eigen, topology, dual-anchor)\n"
     "instead of correcting the vote.",
     "Ranks 24–27 — all BELOW KNNOptK."),
    ("G2", "Magnitude / LocalOdds",  AMBER,
     "Add a secondary distance-magnitude weight or re-estimate\n"
     "the imbalance ratio locally.",
     "Redundant with α-CV in KNNFairRankCV."),
    ("G3", "Density / Bayesian",     AMBER,
     "Estimate local density or impose a prior on r;\n"
     "density estimation noisy in moderate-to-high dim.",
     "All rank below KNNWeighted — noise > signal."),
    ("G4", "Topo / Ensemble redund", GRAY,
     "TopoCount, TopoJoint, JackknifeEnsemble, LocalCount, LID:\n"
     "either dominated by TopoJointBootstrap or aggregate noise.",
     "Strictly dominated on every dataset."),
]
gy = Inches(1.8); gh = Inches(1.15); gw = Inches(12.2)
for i, (idx, name, color, descr, finding) in enumerate(groups_d):
    y = gy + i * (gh + Inches(0.05))
    add_rect(s, Inches(0.55), y, gw, gh, fill=LIGHT)
    add_rect(s, Inches(0.55), y, Inches(1.2), gh, fill=color)
    add_text(s, idx, Inches(0.55), y, Inches(1.2), gh,
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
    add_text(s, name, Inches(1.85), y + Inches(0.12),
             Inches(4), Inches(0.4),
             size=14, bold=True, color=DARK, font=FONT_HEAD)
    add_text(s, descr, Inches(1.85), y + Inches(0.5),
             Inches(7), Inches(0.65),
             size=10, color=GRAY)
    add_text(s, finding, Inches(9), y + Inches(0.12),
             Inches(3.8), Inches(0.95),
             size=11, bold=True, color=NAVY, italic=True,
             valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 27 — ANNEX: Cohen's d full table
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Tabela Cohen's d completa contra cada baseline em cada métrica. "
    "Células a azul escuro: |d| ≥ 0.8 (large). Células a azul claro: "
    "|d| ≥ 0.5 (medium)."
)
slide_header(s, "ANNEX", "Cohen's d  —  effect sizes vs KNNOptK and SMOTE+KNN", 27)

def color_d(v):
    av = abs(v)
    if av >= 0.8:  return NAVY, WHITE
    if av >= 0.5:  return PALE, NAVY
    return WHITE, DARK

d_rows = [
    ("JointCV",         [0.88, 0.34, 0.69], [0.30, 0.32, 0.21]),
    ("Ensemble",        [0.72, 0.56, 0.58], [0.18, 0.56, 0.15]),
    ("CV",              [0.80, 0.32, 0.65], [0.26, 0.30, 0.19]),
    ("OptVotes",        [0.76, 0.30, 0.62], [0.24, 0.27, 0.17]),
    ("FairRank",        [0.62, 0.27, 0.54], [0.15, 0.24, 0.12]),
    ("Jackknife",       [0.61, 0.36, 0.52], [0.13, 0.35, 0.11]),
    ("TopoJB",          [0.58, 0.28, 0.50], [0.12, 0.26, 0.10]),
    ("KNNWeighted",     [0.12, -0.20, 0.05], [-0.45, -0.30, -0.40]),
    ("SMOTE+KNN",       [0.45, 0.05, 0.27],  [None, None, None]),
]
tx, ty = Inches(0.55), Inches(1.7)
total_w = Inches(12.2)
col_w = (total_w - Inches(2.8)) / 6
alg_w = Inches(2.8)
rh_ = Inches(0.42)

add_rect(s, tx, ty, alg_w, rh_ * 2, fill=NAVY)
add_text(s, "Algorithm", tx + Inches(0.2), ty, alg_w, rh_ * 2,
         size=11, bold=True, color=WHITE, valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
add_rect(s, tx + alg_w, ty, col_w * 3, rh_, fill=NAVY)
add_text(s, "vs KNNOptK", tx + alg_w, ty, col_w * 3, rh_,
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
add_rect(s, tx + alg_w + col_w * 3, ty, col_w * 3, rh_, fill=BLUE)
add_text(s, "vs SMOTE+KNN", tx + alg_w + col_w * 3, ty, col_w * 3, rh_,
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
for i, label in enumerate(["G-Mean", "ROC-AUC", "MCC"]):
    add_rect(s, tx + alg_w + i * col_w, ty + rh_, col_w, rh_, fill=NAVY)
    add_text(s, label, tx + alg_w + i * col_w, ty + rh_, col_w, rh_,
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
for i, label in enumerate(["G-Mean", "ROC-AUC", "MCC"]):
    add_rect(s, tx + alg_w + (3 + i) * col_w, ty + rh_, col_w, rh_, fill=BLUE)
    add_text(s, label, tx + alg_w + (3 + i) * col_w, ty + rh_, col_w, rh_,
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
for i, (alg, vs_optk, vs_smote) in enumerate(d_rows):
    y = ty + (i + 2) * rh_
    bg = WHITE if i % 2 == 0 else LIGHT
    add_rect(s, tx, y, alg_w, rh_, fill=bg)
    add_text(s, alg, tx + Inches(0.2), y, alg_w, rh_,
             size=11, bold=True, color=DARK, valign=MSO_ANCHOR.MIDDLE,
             font=FONT_HEAD)
    for j, v in enumerate(vs_optk):
        cell_x = tx + alg_w + j * col_w
        if v is None:
            add_rect(s, cell_x, y, col_w, rh_, fill=bg); continue
        fill, fg = color_d(v)
        add_rect(s, cell_x, y, col_w, rh_, fill=fill)
        add_text(s, f"{v:+.2f}", cell_x, y, col_w, rh_,
                 size=10, bold=(abs(v) >= 0.5), color=fg,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
                 font=FONT_HEAD)
    for j, v in enumerate(vs_smote):
        cell_x = tx + alg_w + (3 + j) * col_w
        if v is None:
            add_rect(s, cell_x, y, col_w, rh_, fill=bg); continue
        fill, fg = color_d(v)
        add_rect(s, cell_x, y, col_w, rh_, fill=fill)
        add_text(s, f"{v:+.2f}", cell_x, y, col_w, rh_,
                 size=10, bold=(abs(v) >= 0.5), color=fg,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
                 font=FONT_HEAD)

add_rect(s, Inches(4), Inches(6.75), Inches(0.4), Inches(0.25), fill=NAVY)
add_text(s, "|d| ≥ 0.8 large",
         Inches(4.5), Inches(6.75), Inches(2), Inches(0.25),
         size=10, color=DARK, valign=MSO_ANCHOR.MIDDLE)
add_rect(s, Inches(6.5), Inches(6.75), Inches(0.4), Inches(0.25), fill=PALE)
add_text(s, "|d| ≥ 0.5 medium",
         Inches(7), Inches(6.75), Inches(2.2), Inches(0.25),
         size=10, color=DARK, valign=MSO_ANCHOR.MIDDLE)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 28 — ANNEX: notebook navigator
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(
    "Para reprodutibilidade: o repositório tem múltiplos notebooks. O "
    "entrega_final.ipynb é o deliverable principal; meta_analysis.ipynb "
    "tem a pipeline ISA; phase2_benchmark.ipynb gera o "
    "benchmark_5rep.csv."
)
slide_header(s, "ANNEX", "Notebook navigator", 28)

nbs = [
    ("entrega_final.ipynb",        "★ Main deliverable — full pipeline + theory + results"),
    ("meta_analysis.ipynb",        "ISA + LODO-CV + stress test + Kendall trend"),
    ("phase2_benchmark.ipynb",     "Benchmark execution → benchmark_5rep.csv"),
    ("phase1_baseline.ipynb",      "Phase 1 baseline analysis (KNNOptK)"),
    ("algorithm_documentation.ipynb", "Per-algorithm technical reference"),
    ("algorithm_design.ipynb",     "Development narrative (legacy)"),
    ("algorithm_reference.ipynb",  "Earlier reference (legacy)"),
]
tx, ty = Inches(0.55), Inches(1.8); rh_ = Inches(0.55); nw = Inches(4.5)
for i, (nb_name, role) in enumerate(nbs):
    y = ty + i * (rh_ + Inches(0.1))
    add_rect(s, tx, y, Inches(12.2), rh_, fill=LIGHT)
    add_rect(s, tx, y, Inches(0.08), rh_, fill=NAVY)
    add_text(s, nb_name, tx + Inches(0.25), y, nw, rh_,
             size=12, bold=True, color=NAVY, valign=MSO_ANCHOR.MIDDLE,
             font=FONT_HEAD)
    add_text(s, role, tx + nw + Inches(0.4), y, Inches(7.4), rh_,
             size=11, color=DARK, valign=MSO_ANCHOR.MIDDLE)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 29 — ANNEX: references
# ═══════════════════════════════════════════════════════════════════════
s = add_slide("Referências principais.")
slide_header(s, "ANNEX", "References", 29)

refs = [
    ("Demšar, J. (2006).",
     "Statistical comparisons of classifiers over multiple data sets.",
     "JMLR 7, 1–30."),
    ("Smith-Miles, K. & Muñoz, M.A. (2023).",
     "Instance Space Analysis: A toolkit for the assessment of algorithmic power.",
     "ACM Computing Surveys 55(12)."),
    ("Hastie, T. & Tibshirani, R. (1996).",
     "Discriminant adaptive nearest-neighbor classification.",
     "IEEE TPAMI 18(6), 607–616."),
    ("Chawla, N.V. et al. (2002).",
     "SMOTE: Synthetic Minority Over-sampling Technique.",
     "JAIR 16, 321–357."),
    ("rushter / MLAlgorithms.",
     "Base KNN loop adapted with sklearn-compatible interface.",
     "github.com/rushter/MLAlgorithms"),
]
ty = Inches(1.85); rh_ = Inches(0.9)
for i, (author, title, venue) in enumerate(refs):
    y = ty + i * (rh_ + Inches(0.1))
    add_rect(s, Inches(0.55), y, Inches(12.2), rh_, fill=LIGHT)
    add_rect(s, Inches(0.55), y, Inches(0.08), rh_, fill=NAVY)
    add_text(s, author, Inches(0.85), y + Inches(0.1),
             Inches(11.5), Inches(0.3),
             size=12, bold=True, color=NAVY, font=FONT_HEAD)
    add_text(s, title, Inches(0.85), y + Inches(0.4),
             Inches(11.5), Inches(0.3),
             size=12, color=DARK, italic=True)
    add_text(s, venue, Inches(0.85), y + Inches(0.65),
             Inches(11.5), Inches(0.3),
             size=11, color=GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 30 — Thank you
# ═══════════════════════════════════════════════════════════════════════
s = add_slide("Encerramento. Obrigado pela atenção. Questões.")
add_rect(s, Inches(0), Inches(0), SW, SH, fill=NAVY)
add_rect(s, Inches(5.665), Inches(2.8), Inches(2.0), Inches(0.08), fill=BLUE)
add_text(s, "Obrigado",
         Inches(0), Inches(2.95), SW, Inches(1.5),
         size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_HEAD)
add_text(s, "Questions ?",
         Inches(0), Inches(4.45), SW, Inches(0.8),
         size=28, color=PALE, italic=True, align=PP_ALIGN.CENTER, font=FONT_HEAD)
add_text(s,
         "Duarte Gomes  ·  José Sousa  ·  Tiago Sousa     |     Group 2  ·  CC2008  ·  2025/2026",
         Inches(0), Inches(6.7), SW, Inches(0.4),
         size=12, color=PALE, align=PP_ALIGN.CENTER, font=FONT_HEAD)


# ───────────────────────────── Save ──────────────────────────────────
prs.save(str(OUTPUT))
print(f"Saved: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
